import os.path
from pptx import Presentation
import PyPDF2
import slate3k as slate
import pandas as pd
from cgitb import text
import nltk
import spacy
import unicodedata
from model.archivo import Archivo
from model.concepto import Concepto
from nltk.corpus import stopwords

def strip_accents(text):

    try:
        text = unicode(text, 'utf-8')
    except NameError: # unicode is a default on python 3 
        pass

    text = unicodedata.normalize('NFD', text)\
           .encode('ascii', 'ignore')\
           .decode("utf-8")

    return str(text)

def process_file(source, db):
    
    sr = stopwords.words('spanish')
    
    #Obtengo referencias al nombre y extensión del archivo
    base=os.path.basename(source)
    nombre=os.path.splitext(base)[0]
    extension=os.path.splitext(base)[1]

    contenido = []

    #Curso en base a la extensión del archivo. 
    if extension=='.pptx':
        prs = Presentation(source)
            
        #Lectura del PPT
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        contenido.append(run.text) 
            contenido.append('.')

        #Elimino residuos del contenido
        indices_a_borrar = []
        for index, item in enumerate(contenido):
            if len(item)<3 and item!='.':
                indices_a_borrar.append(index)

        for index in sorted(indices_a_borrar, reverse=True):
            del contenido[index]

        #Elimino residuos del contenido gral
        contenido= [frase.replace('\n', '') for frase in contenido]
        contenido= [frase.replace('\t', '') for frase in contenido]

        #Genero un solo texto con todo el contenido
        strcont = ' '.join(contenido)
        oraciones = strcont.split('.')

        #Obtenemos los conceptos o palabras que más aparecen en el contenido
        freq = nltk.FreqDist(contenido)

    elif extension=='.pdf':

        pdfFileObj = open(source, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)

        for i in range(0,pdfReader.getNumPages()):
            pageObj = pdfReader.getPage(i)
            contenido.append(pageObj.extractText())

        #Elimino residuos del contenido
        indices_a_borrar = []
        for index, item in enumerate(contenido):
            if len(item)<3 and item!='.':
                indices_a_borrar.append(index)

        for index in sorted(indices_a_borrar, reverse=True):
            del contenido[index]

        #Elimino residuos del contenido gral
        contenido= [frase.replace('\n', '') for frase in contenido]
        contenido= [frase.replace('\t', '') for frase in contenido]

        #Genero un solo texto con todo el contenido
        strcont = ' '.join(contenido)
        oraciones = strcont.split('.')

        words = nltk.tokenize.word_tokenize(strcont)

        clean_tokens = words[:]

        for token in words:
            if token in sr:
                clean_tokens.remove(token)

        freq = nltk.FreqDist(clean_tokens)

    elif extension=='.txt':

        file = open(source,mode='r')
        contenido = file.readlines()
        file.close()

        words=[]

        for linea in contenido:
            words_per_line = linea.split()
            words.append(words_per_line)

        #Elimino residuos del contenido
        indices_a_borrar = []
        for index, item in enumerate(contenido):
            if len(item)<3 and item!='.':
                indices_a_borrar.append(index)

        for index in sorted(indices_a_borrar, reverse=True):
            del contenido[index]

        contenido= [frase.replace('\n', '') for frase in contenido]
        contenido= [frase.replace('\t', '') for frase in contenido]

        strcont = ' '.join(contenido)
        oraciones = strcont.split('.')

        words = nltk.tokenize.word_tokenize(strcont)

        clean_tokens = words[:]

        for token in words:
            if token in sr:
                clean_tokens.remove(token)

        freq = nltk.FreqDist(clean_tokens)

    else:
        print("No se como abrir eso! ")

    #### GENERAL A CUALQUIER ARCHIVO DE CONTENIDO ####

    print("FREQ items: "+str(freq.items()))
        
    claves = []
    topPalabras = '['

    for key,val in sorted(freq.items()):
        if val>1:
            if len(key)>2:
                itemActual='{'+str(key) + ',' + str(val) + '}'
                if itemActual not in sr: 
                    topPalabras+=str(itemActual)
                    claves.append(str(key))

    topPalabras+=']'

    archivo = Archivo(nombre,topPalabras,source)

    db.add_archivo(archivo)
    id_archivo = db.get_last_id()

    for key in claves:
        concepto= key.lower()
        concepto = strip_accents(concepto)
        db.add_concepto(concepto,id_archivo)

        id_concepto = db.get_last_id()

        matching= [s for s in oraciones if key in s]
        for texto in matching:
            texto_a_guardar= texto.lower()
            texto_a_guardar = strip_accents(texto_a_guardar)
            db.add_conceptosxarchivo(id_archivo,id_concepto,texto_a_guardar)

    # The next lines are commented because I did not continue the development of the NLP part

    #nlp=spacy.load('es_core_news_md')
    #doc = nlp(strcont)

    #for token in doc:
       # print(token.text, token.lemma_, token.pos_, token.tag_, token.dep_, token.shape_, token.is_alpha, token.is_stop)

    #print("Verbos:", [token.lemma_ for token in doc if token.pos_ == "VERB"])
    #sentencias = [s for s in doc.sents]